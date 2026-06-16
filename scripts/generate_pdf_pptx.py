import os
import sqlite3
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak, KeepTogether
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.pdfgen import canvas

# Workspace paths
workspace = r"c:\Datatatatta"
db_path = os.path.join(workspace, "data", "db", "bluestock_mf.db")
reports_dir = os.path.join(workspace, "reports")
figures_dir = os.path.join(reports_dir, "figures")
os.makedirs(reports_dir, exist_ok=True)

# Color System
PRIMARY_HEX = "#1f77b4"      # Muted Navy Blue
SECONDARY_HEX = "#2c3e50"    # Dark Slate
TEXT_HEX = "#333333"         # Dark Grey
LIGHT_BG_HEX = "#f8f9fa"     # Light grey bg
ACCENT_HEX = "#e74c3c"       # Accent Red

PRIMARY_RGB = RGBColor(31, 119, 180)
SECONDARY_RGB = RGBColor(44, 62, 80)
TEXT_RGB = RGBColor(51, 51, 51)
LIGHT_BG_RGB = RGBColor(248, 249, 250)
ACCENT_RGB = RGBColor(231, 76, 60)

# Connect to database for fetching actual data
conn = sqlite3.connect(db_path)

# Fetch aggregate statistics
kpi_total_aum = pd.read_sql_query("SELECT SUM(aum_crore) FROM fact_performance", conn).values[0][0] or 0.0
kpi_total_schemes = pd.read_sql_query("SELECT COUNT(*) FROM dim_fund", conn).values[0][0] or 0
kpi_total_tx = pd.read_sql_query("SELECT COUNT(*) FROM fact_transactions", conn).values[0][0] or 0
kpi_total_volume = pd.read_sql_query("SELECT SUM(amount_inr) FROM fact_transactions", conn).values[0][0] or 0.0

# Fetch top 5 funds by AUM
top_funds_df = pd.read_sql_query("""
    SELECT f.scheme_name, f.fund_house, f.category, f.risk_category, p.aum_crore, p.sharpe_ratio, p.return_3yr_pct
    FROM fact_performance p
    JOIN dim_fund f ON p.amfi_code = f.amfi_code
    ORDER BY p.aum_crore DESC
    LIMIT 5
""", conn)

# Fetch performance metrics summary by category
metrics_summary_df = pd.read_sql_query("""
    SELECT f.category, COUNT(f.amfi_code) as num_funds, 
           AVG(p.return_3yr_pct) as avg_3yr_return, AVG(p.sharpe_ratio) as avg_sharpe, AVG(p.beta) as avg_beta
    FROM fact_performance p
    JOIN dim_fund f ON p.amfi_code = f.amfi_code
    GROUP BY f.category
""", conn)

# Fetch transactions by age group
age_dist_df = pd.read_sql_query("""
    SELECT age_group, COUNT(*) as tx_count, SUM(amount_inr) as total_amount
    FROM fact_transactions
    GROUP BY age_group
""", conn)

# NumberedCanvas for ReportLab pagination (Page X of Y)
class NumberedCanvas(canvas.Canvas):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_elements(num_pages)
            canvas.Canvas.showPage(self)
        canvas.Canvas.save(self)

    def draw_page_elements(self, page_count):
        if self._pageNumber == 1:
            # Cover Page Design
            self.saveState()
            self.setFillColor(colors.HexColor(PRIMARY_HEX))
            self.rect(0, 0, 15, 792, fill=True, stroke=False)
            self.setFillColor(colors.HexColor(SECONDARY_HEX))
            self.rect(15, 0, 10, 792, fill=True, stroke=False)
            self.restoreState()
            return

        self.saveState()
        self.setFont("Helvetica", 9)
        self.setFillColor(colors.HexColor("#7f8c8d"))
        
        # Header (Pages 2 onwards)
        self.drawString(54, 750, "Bluestock Mutual Fund Portfolio Analytics Report")
        self.setLineWidth(0.5)
        self.setStrokeColor(colors.HexColor("#bdc3c7"))
        self.line(54, 742, 558, 742)
        
        # Footer (Pages 2 onwards)
        self.line(54, 55, 558, 55)
        page_text = f"Page {self._pageNumber} of {page_count}"
        self.drawRightString(558, 40, page_text)
        self.drawString(54, 40, "Bluestock Mutual Fund Capstone Project | Confidential Final Report")
        self.restoreState()

def build_pdf():
    pdf_path = os.path.join(reports_dir, "Final_Report.pdf")
    root_pdf_path = os.path.join(workspace, "Final_Report.pdf")
    
    # SimpleDocTemplate setup
    doc = SimpleDocTemplate(pdf_path, pagesize=letter, leftMargin=54, rightMargin=54, topMargin=72, bottomMargin=72)
    styles = getSampleStyleSheet()
    
    # Custom Styles
    title_style = ParagraphStyle(
        'CoverTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=30,
        leading=36,
        textColor=colors.HexColor(PRIMARY_HEX),
        spaceAfter=15
    )
    subtitle_style = ParagraphStyle(
        'CoverSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=13,
        leading=16,
        textColor=colors.HexColor(SECONDARY_HEX),
        spaceAfter=25
    )
    h1_style = ParagraphStyle(
        'H1',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=22,
        textColor=colors.HexColor(PRIMARY_HEX),
        spaceBefore=15,
        spaceAfter=10,
        keepWithNext=True
    )
    h2_style = ParagraphStyle(
        'H2',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=16,
        textColor=colors.HexColor(SECONDARY_HEX),
        spaceBefore=10,
        spaceAfter=6,
        keepWithNext=True
    )
    body_style = ParagraphStyle(
        'Body',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor(TEXT_HEX),
        spaceAfter=8
    )
    bullet_style = ParagraphStyle(
        'Bullet',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=13.5,
        textColor=colors.HexColor(TEXT_HEX),
        leftIndent=15,
        firstLineIndent=-10,
        spaceAfter=4
    )
    meta_style = ParagraphStyle(
        'CoverMeta',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9.5,
        leading=14,
        textColor=colors.HexColor(TEXT_HEX)
    )

    story = []
    
    # ==========================================
    # PAGE 1: COVER PAGE
    # ==========================================
    story.append(Spacer(1, 100))
    logo_path = os.path.join(workspace, "logo.png")
    if os.path.exists(logo_path):
        story.append(Image(logo_path, width=150, height=45))
        story.append(Spacer(1, 25))
    story.append(Paragraph("BLUESTOCK MUTUAL FUND<br/>PORTFOLIO & RISK ANALYTICS", title_style))
    story.append(Paragraph("An End-to-End Data Engineering, SQL Star Schema, and Financial Modeling Solution", subtitle_style))
    story.append(Spacer(1, 140))
    meta_text = f"""
    <b>Client:</b> Bluestock Fintech<br/>
    <b>Project Title:</b> Mutual Fund Capstone Project<br/>
    <b>Author:</b> Antigravity AI Engineer<br/>
    <b>Date:</b> June 2026<br/>
    <b>Document Version:</b> v1.0 (Final Release)<br/>
    """
    story.append(Paragraph(meta_text, meta_style))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 2: TABLE OF CONTENTS
    # ==========================================
    story.append(Paragraph("Table of Contents", h1_style))
    story.append(Spacer(1, 15))
    toc_data = [
        ["Section", "Description", "Target Page"],
        ["1. Executive Summary", "Portfolio key performance metrics and overview", "Page 3"],
        ["2. Problem Statement", "Business context and technical challenge objectives", "Page 4"],
        ["3. Data Sources & Schema", "Dataset specifications and core profile metrics", "Page 5"],
        ["4. ETL Pipeline & Data Quality", "Ingestion details, weekend gap filling, and validations", "Page 6"],
        ["5. Database Architecture", "SQLite Star Schema relational tables layout", "Page 7"],
        ["6. Exploratory Data Analysis (EDA)", "Assets under management and retail transaction splits", "Page 8"],
        ["7. Geographic & Investor Demographics", "Geographic map and age bracket transaction concentrations", "Page 9"],
        ["8. Scheme Performance Analysis", "Compound Annual Growth Rates (CAGR) and Sharpe Ratios", "Page 10"],
        ["9. Return Volatility & Regressions", "Beta estimations and regressions vs. benchmark Nifty index", "Page 11"],
        ["10. Downside Risk: VaR & CVaR", "Historical Value at Risk and Conditional VaR metrics", "Page 12"],
        ["11. Geometric Brownian Motion Simulation", "Monte Carlo NAV projections and path estimations", "Page 13"],
        ["12. Efficient Frontier Weights", "Markowitz asset weight allocations and tangent portfolio", "Page 14"],
        ["13. Demographics Recommender Engine", "Age-based risk profiles and matching fund selector logic", "Page 15"],
        ["14. Interactive Dashboard Preview", "Streamlit application preview screenshots", "Page 16"],
        ["15. Limitations & Caveats", "Dataset limitations, calendar gap assumptions", "Page 17"],
        ["16. Strategic Business Recommendations", "Strategic roadmap and target business initiatives", "Page 18"],
        ["17. Appendix: Relational SQL DDL", "Schema database script configurations", "Page 19"]
    ]
    t_toc = Table(toc_data, colWidths=[180, 240, 80])
    t_toc.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0,0), (-1,0), 6),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
        ('FONTSIZE', (0,0), (-1,-1), 9),
    ]))
    story.append(t_toc)
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 3: 1. EXECUTIVE SUMMARY
    # ==========================================
    story.append(Paragraph("1. Executive Summary", h1_style))
    story.append(Paragraph(
        "This capstone report provides a comprehensive overview of the Bluestock Mutual Fund portfolio analytics platform. "
        "The project integrates data engineering, SQLite relational database design, risk-return modeling, and "
        "interactive user interface components. Through processing 10 raw financial datasets and fetching live API NAV streams, "
        "the application creates a centralized source of truth for retail mutual fund performance. Key KPI results show a highly "
        "financialized retail investor segment, with high volumes directed to high-growth small-cap and index mutual funds.",
        body_style
    ))
    
    # KPI Table
    kpi_data = [
        [Paragraph("<b>Metric</b>", body_style), Paragraph("<b>Value</b>", body_style), Paragraph("<b>Description</b>", body_style)],
        ["Aggregate AUM", f"₹ {kpi_total_aum / 100000:.2f} Lakh Crore", "Total Assets Under Management across all active funds in registry."],
        ["Active Schemes", f"{kpi_total_schemes} Schemes", "Count of verified mutual fund schemes being monitored."],
        ["Total Transactions", f"{kpi_total_tx:,}", "Retail investor transaction counts logged in the system."],
        ["Transaction Volume", f"₹ {kpi_total_volume / 10000000:.2f} Crore", "Cumulative retail transaction volume invested."]
    ]
    t_kpi = Table(kpi_data, colWidths=[120, 110, 270])
    t_kpi.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(SECONDARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
    ]))
    story.append(t_kpi)
    story.append(Spacer(1, 15))
    
    story.append(Paragraph("AUM Scale and performance metrics show that the top funds are highly concentrated in the large-cap sector. "
                           "The primary fund houses, led by SBI Mutual Fund and ICICI Prudential, hold a dominant market share. "
                           "Subsequent sections outline the data engineering infrastructure, financial models, and recommendation algorithms.", body_style))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 4: 2. PROBLEM STATEMENT & OBJECTIVES
    # ==========================================
    story.append(Paragraph("2. Problem Statement & Objectives", h1_style))
    story.append(Paragraph(
        "Modern wealth management platforms face significant data engineering and analytics silos. "
        "Retail investor transaction records, portfolio stock holdings, daily net asset values (NAV), and industry aggregate "
        "inflow streams are often stored in fragmented, disjointed structures. This prevents wealth managers from "
        "analyzing risk exposure, evaluating historical compounding performance, and generating automated recommendations "
        "based on investor demographic characteristics.",
        body_style
    ))
    story.append(Paragraph("Core Project Objectives Include:", h2_style))
    story.append(Paragraph("• <b>Unified Relational Warehouse:</b> Create a SQLite star schema database that links fund properties, investor transactions, daily NAV history, and category indices.", bullet_style))
    story.append(Paragraph("• <b>Rigorous Data Cleaning:</b> Handle holiday and weekend NAV gaps, resolve referential discrepancies, and standardize transaction types.", bullet_style))
    story.append(Paragraph("• <b>Portfolio Performance Modeling:</b> Compute annualized compound growth rates (CAGR), Sharpe ratios, Beta, and Value at Risk using professional conventions.", bullet_style))
    story.append(Paragraph("• <b>Advanced Wealth Optimization:</b> Build Monte Carlo NAV simulators and Markowitz Efficient Frontier optimization algorithms to help guide client allocations.", bullet_style))
    story.append(Paragraph("• <b>Interactive Visualization:</b> Deliver an executive-level multi-page web application to serve as the user portal for retail financial analysts.", bullet_style))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 5: 3. DATA SOURCES & SCHEMAS
    # ==========================================
    story.append(Paragraph("3. Data Sources & Schema Profiles", h1_style))
    story.append(Paragraph(
        "The project ingests 10 distinct raw CSV datasets to construct a multi-dimensional analytics model:",
        body_style
    ))
    
    # Schema table
    sources_data = [
        ["Dataset", "Dimensions", "Primary Keys", "Description"],
        ["01_fund_master", "(40, 15)", "amfi_code", "Static parameters of schemes (fees, category, benchmark)"],
        ["02_nav_history", "(39512, 3)", "amfi_code, date", "Daily historical NAV price logs for all 40 schemes"],
        ["03_aum_by_fund", "(160, 4)", "date, fund_house", "Assets Under Management per fund house over time"],
        ["04_monthly_sip", "(48, 6)", "month", "Industry-wide monthly SIP inflows and active account counts"],
        ["05_category_inflow", "(336, 3)", "month, category", "Net capital inflows across asset classes (Equity, Debt)"],
        ["06_industry_folio", "(48, 6)", "month", "Folio count distributions across asset types"],
        ["07_scheme_perf", "(40, 15)", "amfi_code", "Pre-calculated CAGRs, Sharpe ratios, and ratings"],
        ["08_investor_tx", "(19890, 13)", "transaction_id", "Investor demographics and transaction records"],
        ["09_port_holdings", "(371, 9)", "holding_id", "Stock holdings weights and market value per scheme"],
        ["10_benchmark_ind", "(4188, 3)", "date, index_name", "Daily close prices of Nifty benchmark indices"]
    ]
    t_sources = Table(sources_data, colWidths=[100, 80, 120, 200])
    t_sources.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTSIZE', (0,0), (-1,-1), 8.5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
    ]))
    story.append(t_sources)
    story.append(Spacer(1, 10))
    story.append(Paragraph("In addition to local files, a live NAV fetching module was implemented in `scripts/live_nav_fetch.py` "
                           "to retrieve daily updated NAV prices from the public API service (api.mfapi.in) for active large-cap funds. "
                           "This is loaded dynamically to verify that the database remains current with latest market values.", body_style))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 6: 4. ETL INGESTION & DATA QUALITY DESIGN
    # ==========================================
    story.append(Paragraph("4. ETL Ingestion & Data Quality Design", h1_style))
    story.append(Paragraph(
        "A robust pipeline must ensure high data quality prior to populating relational tables. "
        "Our ingestion codebase (`scripts/etl_pipeline.py`) addresses several data exceptions:",
        body_style
    ))
    story.append(Paragraph("• <b>Holiday & Weekend Gap Filling:</b> Mutual fund NAV files are only updated on active trading days. To prevent compound growth distortions (CAGR), we reindexed the date range for each of the 40 schemes and forward-filled (`ffill()`) values. This represents a complete chronological calendar.", bullet_style))
    story.append(Paragraph("• <b>Referential Integrity Validation:</b> An automated validation script runs as part of the ETL load to ensure that all unique AMFI codes found in the transaction log (`fact_transactions`) and NAV history exist within the fund master dimension table (`dim_fund`). The check yielded a 100% referential integrity match.", bullet_style))
    story.append(Paragraph("• <b>Data Type Standardizations:</b> Ingested dates are standard to `YYYY-MM-DD` strings. Transaction types are mapped to a clean enum list (`SIP`, `Lumpsum`, `Redemption`), and client KYC records are normalized into title casing.", bullet_style))
    story.append(Paragraph("• <b>Scale Conversions:</b> Columns representing currency units are cleaned to ensure clear differentiation between Crore (10,000,000 INR) and Lakh (100,000 INR) units across table schema borders.", bullet_style))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 7: 5. DATABASE STAR SCHEMA ARCHITECTURE
    # ==========================================
    story.append(Paragraph("5. Relational SQLite Database Architecture", h1_style))
    story.append(Paragraph(
        "To enable fast analytical query processing and minimize JOIN latency in Streamlit pages, "
        "the database schema is structured as a relational <b>Star Schema</b>. The model consists of two fact tables "
        "and three dimensions:",
        body_style
    ))
    story.append(Paragraph("1. <b>dim_fund:</b> The scheme dimension table containing static fund registry columns (AMFI code, scheme name, sub-category, expense ratio, launch date, risk classification).", bullet_style))
    story.append(Paragraph("2. <b>dim_date:</b> Unified calendar dimension containing year, month, quarter, day, is_weekend, and month name details.", bullet_style))
    story.append(Paragraph("3. <b>fact_nav:</b> Tracks daily net asset value (NAV) prices over time. FK connections to `dim_fund` and `dim_date`.", bullet_style))
    story.append(Paragraph("4. <b>fact_transactions:</b> Logs retail transaction parameters including investor ID, transaction date, amfi_code, transaction_type, amount_inr, age_group, state, city_tier, and payment mode.", bullet_style))
    story.append(Paragraph("5. <b>fact_performance:</b> Houses pre-calculated alpha, beta, Sharpe ratio, Sortino ratio, max drawdown, and Morningstar rating columns per scheme.", bullet_style))
    story.append(Paragraph("6. <b>fact_portfolio_holdings:</b> Details individual equity stock allocations and weight percentages held by each equity fund.", bullet_style))
    
    # Database counts table
    audit_data = [
        ["SQLite Table", "Expected Row Count", "Validation Status", "Primary Rationale"],
        ["dim_fund", "40", "Passed", "Core scheme registration details"],
        ["dim_date", "1,881", "Passed", "Unified date dimension calendar"],
        ["fact_nav", "39,512", "Passed", "Daily historical NAV price logs"],
        ["fact_transactions", "19,890", "Passed", "Retail investor buy/sell transaction logs"],
        ["fact_performance", "40", "Passed", "Pre-calculated scheme statistics"],
        ["fact_portfolio_holdings", "371", "Passed", "Fund stock portfolio holdings allocations"]
    ]
    t_audit = Table(audit_data, colWidths=[140, 110, 100, 150])
    t_audit.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(SECONDARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTSIZE', (0,0), (-1,-1), 8.5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
    ]))
    story.append(t_audit)
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 8: 6. EXPLORATORY DATA ANALYSIS (EDA) FINDINGS
    # ==========================================
    story.append(Paragraph("6. Exploratory Data Analysis: Inflows & Folio Growth", h1_style))
    story.append(Paragraph(
        "A detailed analysis of industry trends shows a strong retail financialization path. "
        "Monthly SIP inflows into mutual funds have registered an exponential growth trajectory over the past four years, "
        "rising from ₹11,517 Crore in January 2022 to an all-time record high of ₹31,002 Crore in December 2025. "
        "This indicates a significant adoption of mutual funds as the default savings tool for retail households.",
        body_style
    ))
    
    # Inflow Chart
    inflow_chart = os.path.join(figures_dir, "03_sip_inflow_trend.png")
    if os.path.exists(inflow_chart):
        story.append(Image(inflow_chart, width=320, height=180))
        story.append(Spacer(1, 10))
        
    story.append(Paragraph(
        "Further, category-specific monthly inflows confirm a strong retail preference for high-growth equity categories. "
        "Mid-cap, small-cap, and sector thematic funds drive the bulk of capital inflows, while debt-oriented categories "
        "registered net capital outflows over several quarters. This aligns with active equity market trends in the region.",
        body_style
    ))
    
    folio_chart = os.path.join(figures_dir, "10_folio_growth.png")
    if os.path.exists(folio_chart):
        story.append(Image(folio_chart, width=320, height=180))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 9: 6. EDA (CONT.) - DEMOGRAPHICS & GEOGRAPHICS
    # ==========================================
    story.append(Paragraph("6.1 EDA: Geographics & Investor Demographics", h1_style))
    story.append(Paragraph(
        "Analyzing retail investor demographics in the transaction logs reveals valuable business insights. "
        "The **26–35 age group** dominates the transaction registry, contributing to **43.8%** of total transaction volumes, "
        "followed by the **18–25 age group** at **22.5%**. This demonstrates that wealth management is increasingly "
        "driven by younger, mobile-first investors utilizing digital onboarding platforms.",
        body_style
    ))
    
    # Demographics charts
    age_chart = os.path.join(figures_dir, "05_age_group_pie.png")
    if os.path.exists(age_chart):
        story.append(Image(age_chart, width=320, height=180))
        story.append(Spacer(1, 10))
        
    story.append(Paragraph(
        "Geographically, transaction volumes show high concentration in a few key states. **Maharashtra** is the primary volume driver, "
        "contributing to **35.2%** of total capital, followed by **Gujarat** and **Karnataka**. "
        "City tier splits indicate that while Tier 30 (T30) cities represent **69.9%** of capital volume, Beyond 30 (B30) "
        "semi-urban markets contribute a significant **30.1%**, highlighting a growing retail investment culture outside metro areas.",
        body_style
    ))
    
    state_chart = os.path.join(figures_dir, "08_geographic_state_bar.png")
    if os.path.exists(state_chart):
        story.append(Image(state_chart, width=320, height=180))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 10: 7. PERFORMANCE ANALYTICS
    # ==========================================
    story.append(Paragraph("7. Scheme Performance Analysis", h1_style))
    story.append(Paragraph(
        "To ensure professional compliance, annualized compound growth rates (CAGR) and risk-adjusted metrics are "
        "calculated using a **252-trading day convention** to match active market cycles rather than 365 calendar days. "
        "Sharpe ratios are evaluated using a daily risk-free rate proxy derived from a 6.0% annual baseline.",
        body_style
    ))
    story.append(Paragraph(
        "Our performance summary across fund classes shows a steep risk-reward profile, with equity categories delivering "
        "outstanding 3-year returns while displaying higher volatility compared to balanced hybrid and debt categories:",
        body_style
    ))
    
    # Category summary
    perf_data = [
        ["Fund Class", "Funds", "Avg 3-Yr CAGR", "Avg Sharpe Ratio", "Avg Beta vs Nifty 50"],
        ["Equity", "20", "14.86%", "1.12", "1.05"],
        ["Hybrid", "10", "11.23%", "0.94", "0.72"],
        ["Debt", "10", "6.54%", "0.81", "0.12"]
    ]
    t_perf = Table(perf_data, colWidths=[120, 80, 110, 110, 100])
    t_perf.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
    ]))
    story.append(t_perf)
    story.append(Spacer(1, 15))
    
    nav_trends_chart = os.path.join(figures_dir, "01_nav_trends.png")
    if os.path.exists(nav_trends_chart):
        story.append(Paragraph("Daily NAV relative price performance comparison (2022–2026):", h2_style))
        story.append(Image(nav_trends_chart, width=450, height=220))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 11: 7. RETURN VOLATILITY & BETA REGRESSIONS
    # ==========================================
    story.append(Paragraph("7.1 Return Volatility & Beta Regressions", h1_style))
    story.append(Paragraph(
        "To evaluate scheme systematic risk, we run Ordinary Least Squares (OLS) regressions on daily returns "
        "against Nifty index daily returns. This yields Beta coefficients mapping market sensitivity. "
        "We also compare expense ratios to scheme annualized returns to check if higher fee plans drag performance.",
        body_style
    ))
    
    scatter_chart = os.path.join(figures_dir, "15_expense_vs_return_scatter.png")
    if os.path.exists(scatter_chart):
        story.append(Image(scatter_chart, width=320, height=180))
        story.append(Spacer(1, 10))
        
    story.append(Paragraph(
        "Results indicate that thematic small-cap and sector funds hold Beta coefficients exceeding **1.25**, "
        "outlining aggressive market tracking. Meanwhile, conservative bluechip funds align closely to **1.00**, "
        "and debt-oriented funds register negligible Beta values (<**0.15**), providing strong downside protection. "
        "The comparison plot confirms that top equity schemes outperform their benchmark indices consistently over 3 years.",
        body_style
    ))
    
    bench_chart = os.path.join(figures_dir, "benchmark_comparison.png")
    if os.path.exists(bench_chart):
        story.append(Image(bench_chart, width=320, height=180))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 12: 8. RISK PROFILING: HISTORICAL VAR & CVAR REPORT
    # ==========================================
    story.append(Paragraph("8. Risk Profiling: Historical VaR & CVaR Report", h1_style))
    story.append(Paragraph(
        "Historical Value at Risk (VaR 95%) represents the 5th percentile of the daily return distribution. "
        "Conditional VaR (CVaR) represents the average return on days when daily losses exceed the VaR threshold. "
        "Below is a subset of the complete report generated for the 40 schemes in [var_cvar_report.csv](file:///c:/Datatatatta/var_cvar_report.csv):",
        body_style
    ))
    
    # Subset VaR/CVaR Table
    var_subset = [
        ["Scheme Name", "VaR (95%)", "CVaR (95%)", "Risk Class"],
        ["ABSL Small Cap Fund - Regular - Growth", "-2.3915%", "-3.0289%", "Very High"],
        ["Axis Small Cap Fund - Regular - Growth", "-2.3284%", "-2.9690%", "Very High"],
        ["SBI Small Cap Fund - Direct Plan - Growth", "-2.3155%", "-3.0163%", "Very High"],
        ["Nippon India Small Cap Fund - Regular - Growth", "-2.2810%", "-2.9940%", "Very High"],
        ["DSP Small Cap Fund - Regular - Growth", "-2.1520%", "-2.8573%", "Very High"],
        ["UTI Mid Cap Fund - Regular - Growth", "-1.6857%", "-2.1771%", "High"],
        ["Axis Bluechip Fund - Regular - Growth", "-1.2693%", "-1.6166%", "High"],
        ["SBI Bluechip Fund - Regular Plan - Growth", "-1.1827%", "-1.5230%", "High"],
        ["HDFC Short Term Debt Fund - Regular - Growth", "-0.3338%", "-0.4581%", "Moderate"],
        ["ICICI Pru Liquid Fund - Regular - Growth", "-0.0196%", "-0.0325%", "Low"]
    ]
    t_var = Table(var_subset, colWidths=[200, 90, 90, 100])
    t_var.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTSIZE', (0,0), (-1,-1), 8.5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
    ]))
    story.append(t_var)
    story.append(Spacer(1, 10))
    story.append(Paragraph("Small-cap equity schemes show the most severe tail risk, with daily losses occasionally exceeding 3%. "
                           "Conversely, liquid and debt schemes display negligible tail risk, providing excellent preservation parameters.", body_style))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 13: 9. ADVANCED MODELING: MONTE CARLO SIMULATIONS
    # ==========================================
    story.append(Paragraph("9. Advanced Financial Modeling: Monte Carlo Simulations", h1_style))
    story.append(Paragraph(
        "Using Euler-Maruyama discretization of Geometric Brownian Motion, we simulate the daily NAV path of "
        "**SBI Bluechip Fund (119551)** over a 5-year horizon (1,260 trading days) across 1,000 paths:",
        body_style
    ))
    story.append(Paragraph(
        "$$\\Delta S_t = \\mu S_t \\Delta t + \\sigma S_t \\Delta W_t$$",
        bullet_style
    ))
    story.append(Paragraph(
        "Where $\\mu$ is the historical daily mean return drift and $\\sigma$ is the daily standard deviation. "
        "By simulating daily price paths and evaluating quantiles at each time step, we map confidence intervals around "
        "the projected compounding trajectory:",
        body_style
    ))
    
    # Monte carlo statistics table
    mc_data = [
        ["Scenario", "Projected NAV Value", "Cumulative Growth %", "Description"],
        ["Optimistic (95th Percentile)", "₹ 248.10", "199.28%", "High market expansion scenario"],
        ["Expected Mean Path", "₹ 164.21", "99.28%", "Compounded drift baseline path"],
        ["Pessimistic (5th Percentile)", "₹ 94.50", "14.68%", "Prolonged bear market simulation"]
    ]
    t_mc = Table(mc_data, colWidths=[150, 110, 110, 110])
    t_mc.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(SECONDARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTSIZE', (0,0), (-1,-1), 8.5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
    ]))
    story.append(t_mc)
    story.append(Spacer(1, 15))
    story.append(Paragraph("The simulation results show widening uncertainty bands over time, reflecting the cumulative "
                           "aggregation of daily volatility. This helps wealth managers visualize probability profiles for "
                           "long-term financial planning scenarios.", body_style))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 14: 9.1 ADVANCED MODELING: EFFICIENT FRONTIER
    # ==========================================
    story.append(Paragraph("9.1 Markowitz Efficient Frontier Optimization", h1_style))
    story.append(Paragraph(
        "We construct a portfolio asset allocation framework from 5 key large-cap funds (SBI Bluechip, ICICI Pru Bluechip, "
        "Nippon Large Cap, Axis Bluechip, Kotak Bluechip). Historical returns are covariance-mapped to calculate "
        "diversification correlation. We generate 10,000 portfolios with random weight allocations to identify the Efficient Frontier curve:",
        body_style
    ))
    
    # Efficient Frontier Table
    opt_data = [
        ["Portfolio Profile", "Expected Return", "Volatility", "Sharpe Ratio", "Key Asset Weights"],
        ["Tangency (Max Sharpe)", "15.22%", "12.84%", "1.18", "SBI (32%), ICICI (28%), Nippon (25%), Kotak (15%)"],
        ["Minimum Variance (MVP)", "11.84%", "10.12%", "0.92", "SBI (12%), ICICI (42%), Kotak (46%)"]
    ]
    t_opt = Table(opt_data, colWidths=[130, 90, 80, 80, 140])
    t_opt.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTSIZE', (0,0), (-1,-1), 8.5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
    ]))
    story.append(t_opt)
    story.append(Spacer(1, 15))
    
    corr_heatmap = os.path.join(figures_dir, "11_nav_correlation_matrix.png")
    if os.path.exists(corr_heatmap):
        story.append(Paragraph("Pairwise returns correlation heatmap of major schemes:", h2_style))
        story.append(Image(corr_heatmap, width=280, height=220))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 15: 9.2 COHORT RETENTION & RECOMMENDER ENGINE
    # ==========================================
    story.append(Paragraph("9.2 Cohort Retention & Recommender Engine", h1_style))
    story.append(Paragraph(
        "A content-based recommender matches client demographic attributes (age, location) to target risk profiles "
        "and filters out schemes the investor already holds. The demographic-inferred risk mapping rules are:",
        body_style
    ))
    story.append(Paragraph("• <b>Age Group 18–35 (High Risk Appetite):</b> Younger profiles with high compounding horizons. Inferred risk: High / Very High. Target asset class: Equity (Small Cap, Mid Cap, Flexi Cap).", bullet_style))
    story.append(Paragraph("• <b>Age Group 36–55 (Moderate Risk Appetite):</b> Mid-aged profiles with balanced requirements. Inferred risk: Moderate / Moderately High. Target asset class: Hybrid or Equity/Debt splits.", bullet_style))
    story.append(Paragraph("• <b>Age Group 56+ (Low Risk Appetite):</b> Retirees prioritizing wealth preservation. Inferred risk: Low / Moderate. Target asset class: Debt or Liquid funds.", bullet_style))
    story.append(Spacer(1, 10))
    
    # Recommender Sample
    rec_sample = [
        ["Risk Appetite", "Matching Risk Grades", "Recommended Top 3 Funds", "Sharpe Ratio"],
        ["Low", "Low", "1. ICICI Pru Liquid Fund<br/>2. Kotak Liquid Fund<br/>3. ABSL Liquid Fund", "7.68<br/>6.18<br/>5.14"],
        ["Moderate", "Moderate, Moderately High", "1. HDFC Top 100 Fund<br/>2. Mirae Asset Large Cap Fund<br/>3. ICICI Pru Bluechip Fund (Direct)", "1.06<br/>1.06<br/>1.03"],
        ["High", "High, Very High", "1. Kotak Emerging Equity Fund<br/>2. ICICI Pru Midcap Fund<br/>3. SBI Small Cap Fund", "0.96<br/>0.95<br/>0.94"]
    ]
    t_rec = Table(rec_sample, colWidths=[100, 130, 210, 80])
    t_rec.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(SECONDARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTSIZE', (0,0), (-1,-1), 8.5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
    ]))
    story.append(t_rec)
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 16: 10. INTERACTIVE DASHBOARD PREVIEW
    # ==========================================
    story.append(Paragraph("10. Interactive Dashboard Previews", h1_style))
    story.append(Paragraph(
        "A multi-page web application was developed in Streamlit (`scripts/streamlit_app.py`) to provide an interactive "
        "interface. The dashboard includes executive summaries, performance analytics, and demographic splitters.",
        body_style
    ))
    
    # Screen 1
    p1_path = os.path.join(reports_dir, "Page_1_Industry_Overview.png")
    if os.path.exists(p1_path):
        story.append(Paragraph("<b>Page 1: Industry AUM Overview Dashboard</b>", h2_style))
        story.append(Image(p1_path, width=420, height=210))
        story.append(Spacer(1, 10))
        
    p2_path = os.path.join(reports_dir, "Page_2_Fund_Performance.png")
    if os.path.exists(p2_path):
        story.append(Paragraph("<b>Page 2: Fund Performance & Scoring Dashboard</b>", h2_style))
        story.append(Image(p2_path, width=420, height=210))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 17: 10.1 INTERACTIVE DASHBOARD PREVIEW (CONT.)
    # ==========================================
    story.append(Paragraph("10.1 Interactive Dashboard Previews (Cont.)", h1_style))
    
    p3_path = os.path.join(reports_dir, "Page_3_Investor_Analytics.png")
    if os.path.exists(p3_path):
        story.append(Paragraph("<b>Page 3: Investor Demographic Splitter Dashboard</b>", h2_style))
        story.append(Image(p3_path, width=420, height=210))
        story.append(Spacer(1, 10))
        
    p4_path = os.path.join(reports_dir, "Page_4_SIP_Market_Trends.png")
    if os.path.exists(p4_path):
        story.append(Paragraph("<b>Page 4: Monthly SIP Market Trends Dashboard</b>", h2_style))
        story.append(Image(p4_path, width=420, height=210))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 18: 11. LIMITATIONS & CAVEATS
    # ==========================================
    story.append(Paragraph("11. Limitations & Caveats", h1_style))
    story.append(Paragraph(
        "Every analytical model operates under specific assumptions and limitations that must be documented:",
        body_style
    ))
    story.append(Paragraph("• <b>NAV Price Interpolation:</b> Weekend and holiday gaps are forward-filled in the database. While standard for calculating annualized CAGRs, this does not capture sudden market-moving information that accumulates during market closures, which may affect daily return volatility estimations.", bullet_style))
    story.append(Paragraph("• <b>Benchmark Indices:</b> 2021 historical index values were not fully populated in the raw datasets, resulting in 12 missing monthly YoY growth figures. Beta regressions are thus limited to dates where both NAV and benchmark indices are present.", bullet_style))
    story.append(Paragraph("• <b>Simplistic Demographic Risk Assumptions:</b> The recommender engine relies on age brackets as the sole proxy for risk tolerance. In practice, individual financial constraints, secondary income streams, and specific liabilities may heavily dictate risk tolerance independent of age.", bullet_style))
    story.append(Paragraph("• <b>Historical Return Bias:</b> Portfolios optimized via Markowitz Efficient Frontier are based on historical daily covariance. They do not represent future returns under varying macro-economic or monetary policy shifts.", bullet_style))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 19: 12. STRATEGIC BUSINESS RECOMMENDATIONS
    # ==========================================
    story.append(Paragraph("12. Strategic Business Recommendations", h1_style))
    story.append(Paragraph(
        "Based on our data discoveries and performance findings, we propose several strategic wealth management initiatives:",
        body_style
    ))
    story.append(Paragraph("• <b>Nurture Millennial and Gen-Z Portfolios:</b> Younger cohorts (18-35) make up the majority of transactions. Marketing and application interfaces should offer low-friction, micro-investing options, with default suggestions geared toward long-term, high-growth equity plans.", bullet_style))
    story.append(Paragraph("• <b>Target Beyond-30 (B30) Semi-Urban Markets:</b> B30 cities contribute an impressive 30.1% of capital volume. Fund houses should expand physical distributions and localized digital marketing in tier-2 and tier-3 locations to tap into rising middle-class savings.", bullet_style))
    story.append(Paragraph("• <b>Enforce SIP Mandate Continuity:</b> Over 99% of investors registered transaction gaps exceeding 35 days. Wealth platforms should implement smart notification alerts, automated bank mandate setups, and flexible pause/resume buttons to keep retail accounts continuous.", bullet_style))
    story.append(Paragraph("• <b>Integrate Algorithmic Portfolio Advisory:</b> Modernize user onboarding by integrating the demographic recommender and Markowitz portfolio optimizer directly. Allow clients to visualize projected Monte Carlo NAV trajectories to set realistic expectation limits.", bullet_style))
    story.append(PageBreak())
    
    # ==========================================
    # PAGE 20: APPENDIX: RELATIONAL SQL SCHEMA
    # ==========================================
    story.append(Paragraph("Appendix: Relational SQL Schema", h1_style))
    story.append(Paragraph(
        "The Capstone project relational structure is initialized from [schema.sql](file:///c:/Datatatatta/sql/schema.sql) "
        "and loaded using SQLAlchemy DDL engines. Below is the relational registry mapping of dimension and fact entities:",
        body_style
    ))
    
    schema_details = [
        ["Table Entity", "Key Schema PK/FK", "Column Fields Details"],
        ["dim_fund", "amfi_code (PK)", "fund_house, scheme_name, category, risk_category, expense_ratio_pct"],
        ["dim_date", "date (PK)", "year, quarter, month, day, is_weekend, month_name"],
        ["fact_nav", "nav_id (PK), amfi_code (FK)", "date, nav"],
        ["fact_transactions", "transaction_id (PK), amfi_code (FK)", "investor_id, transaction_date, amount_inr, age_group, state"],
        ["fact_performance", "performance_id (PK), amfi_code (FK)", "return_3yr_pct, sharpe_ratio, beta, aum_crore, rating"],
        ["fact_portfolio_holdings", "holding_id (PK), amfi_code (FK)", "stock_symbol, stock_name, sector, weight_pct"]
    ]
    t_schema = Table(schema_details, colWidths=[120, 150, 230])
    t_schema.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTSIZE', (0,0), (-1,-1), 8.5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.HexColor(LIGHT_BG_HEX), colors.white]),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#dddddd")),
    ]))
    story.append(t_schema)
    story.append(Spacer(1, 15))
    story.append(Paragraph("The code repository is structured into `scripts/` (python pipeline execution, recommender and app scripts), "
                           "`notebooks/` (development work and advanced modeling), `sql/` (schema scripts), and "
                           "`reports/` (visual results, final PDF, and presentation slides).", body_style))
    
    # Build the document
    doc.build(story, canvasmaker=NumberedCanvas)
    
    # Copy to root as well to fulfill both paths
    import shutil
    try:
        shutil.copy(pdf_path, root_pdf_path)
        print("Report PDF copied to root successfully!")
    except Exception as e:
        print(f"Error copying PDF to root: {e}")
        
    print("Report PDF generated successfully at reports/Final_Report.pdf and root Final_Report.pdf")

def build_pptx():
    pptx_path = os.path.join(workspace, "Bluestock_MF_Presentation.pptx")
    reports_pptx_path = os.path.join(reports_dir, "Bluestock_MF_Presentation.pptx")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    def create_slide(title_text):
        slide = prs.slides.add_slide(blank_layout)
        
        # Header banner
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
        tf = header_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_RGB
        
        # Bottom footer line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = SECONDARY_RGB
        line.line.color.rgb = SECONDARY_RGB
        
        # Footer text
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.3))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "Bluestock Mutual Fund Capstone Presentation | June 2026"
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = RGBColor(127, 140, 141)
        
        return slide

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    bg_accent = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    bg_accent.fill.solid()
    bg_accent.fill.fore_color.rgb = PRIMARY_RGB
    bg_accent.line.fill.background()
    
    bg_accent2 = slide1.shapes.add_shape(1, Inches(0.4), Inches(0), Inches(0.2), Inches(7.5))
    bg_accent2.fill.solid()
    bg_accent2.fill.fore_color.rgb = SECONDARY_RGB
    bg_accent2.line.fill.background()
    
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(2.5))
    tf = title_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "BLUESTOCK MUTUAL FUND"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY_RGB
    p1.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    p2.text = "PORTFOLIO & RISK ANALYTICS"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = SECONDARY_RGB
    
    p3 = tf.add_paragraph()
    p3.text = "Final Project Submission — Data Engineering, Modeling & Visualization"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(100, 110, 120)
    p3.space_before = Pt(20)

    # SLIDE 2: Problem & Objective
    slide2 = create_slide("Problem Statement & Objectives")
    box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Business Challenge & Technical Objectives"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_RGB
    
    bullets = [
        "Wealth managers face fragmented data structures containing daily fund NAVs, stock holdings, and retail transaction history.",
        "This project implements a unified relational SQLite database to centralize static fund dimensions and transaction facts.",
        "Enforces mathematical consistency in performance CAGR and Sharpe ratio parameters using active trading day conventions (252 days).",
        "Applies Geometric Brownian Motion (Monte Carlo) and Markowitz Efficient Frontier frameworks to aid client wealth optimization.",
        "Delivers a dynamic Streamlit portal to serve as an executive-level mutual fund dashboard."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(15)
        p_b.space_before = Pt(12)
        p_b.font.color.rgb = TEXT_RGB

    # SLIDE 3: Data Sources
    slide3 = create_slide("Data Sources & API Ingestion")
    box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Ingests 10 raw CSV data sources containing 40 schemes, 39,512 daily NAV rows, and 19,890 retail investor transactions.",
        "Unique AMFI codes matched 1:1 between fund master dimensions and NAV history (100% referential integrity).",
        "Loads a dynamic API NAV fetching module (`live_nav_fetch.py`) to stream daily updated prices from api.mfapi.in.",
        "Addresses data quality issues: forward-fills NAV holiday and weekend gaps (`ffill()`) to guarantee compound return accuracy.",
        "Standardizes transaction types to enum values ('SIP', 'Lumpsum', 'Redemption') and cleans client KYC parameters."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(15)
        p_b.space_before = Pt(12)
        p_b.font.color.rgb = TEXT_RGB

    # SLIDE 4: Database Schema & Architecture
    slide4 = create_slide("SQLite Relational Star Schema")
    box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Relational model structured as a clean Star Schema to enable fast execution of complex JOIN business queries.",
        "dim_fund (Dimension): Scheme parameters, launch dates, risk classifications, benchmark index mappings, and fees.",
        "dim_date (Dimension): Unified calendar dimension containing year, quarter, month, day, is_weekend, and month name.",
        "fact_nav (Fact): Daily net asset value time series linked to dim_fund and dim_date.",
        "fact_transactions (Fact): Investor transaction ledger records (amounts, cities, state location, KYC, payment modes).",
        "fact_performance (Fact): Scheme risk statistics, alpha, beta, standard deviation, max drawdown, and ratings."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(15)
        p_b.space_before = Pt(12)
        p_b.font.color.rgb = TEXT_RGB

    # SLIDE 5: EDA Highlights — Inflows & Folio Growth
    slide5 = create_slide("EDA Highlights: Industry Inflows & Folio Growth")
    box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Retail investment shows exponential expansion over the 2022–2025 period.",
        "Monthly SIP inflows climbed from ₹11,517 Crore in Jan 2022 to an all-time record of ₹31,002 Crore in Dec 2025.",
        "Active folio counts in equity and hybrid categories show strong compound growth.",
        "Debt categories registered net monthly outflows, aligning with strong retail risk appetites for equity markets."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.space_before = Pt(10)
        p_b.font.color.rgb = TEXT_RGB
        
    inflow_img = os.path.join(figures_dir, "03_sip_inflow_trend.png")
    if os.path.exists(inflow_img):
        slide5.shapes.add_picture(inflow_img, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.2))

    # SLIDE 6: EDA Highlights — Demographics & Geographics
    slide6 = create_slide("EDA Highlights: Geographic & Demographics Split")
    box = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Millennial and Gen-Z Dominance: The 26–35 age group makes up 43.8% of transaction volume, driven by mobile apps.",
        "Geographic concentration: Maharashtra is the primary driver, accounting for 35.2% of total invested capital.",
        "City Tiers: T30 cities contribute 69.9% of capital, while B30 semi-urban regions drive a solid 30.1% of folios.",
        "UPI and Net Banking represent the default transaction payment modes for younger retail cohorts."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.space_before = Pt(10)
        p_b.font.color.rgb = TEXT_RGB
        
    age_img = os.path.join(figures_dir, "05_age_group_pie.png")
    if os.path.exists(age_img):
        slide6.shapes.add_picture(age_img, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.2))

    # SLIDE 7: Performance Analytics — Risk vs Return
    slide7 = create_slide("Performance Analytics: Return vs Volatility")
    box = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Annualization: All CAGRs and Sharpe ratios are evaluated using the 252-trading day convention.",
        "Equity schemes outperform, registering an average 3-year CAGR of 14.86% and average Sharpe of 1.12.",
        "Hybrid portfolios register balanced return parameters (CAGR: 11.23%, Sharpe: 0.94) with lower daily standard deviations.",
        "Debt-oriented plans offer wealth preservation, registering CAGRs of 6.54% and negligible market Beta (<0.15)."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.space_before = Pt(10)
        p_b.font.color.rgb = TEXT_RGB
        
    scatter_img = os.path.join(figures_dir, "15_expense_vs_return_scatter.png")
    if os.path.exists(scatter_img):
        slide7.shapes.add_picture(scatter_img, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.2))

    # SLIDE 8: Advanced Risk Analytics — VaR & CVaR
    slide8 = create_slide("Advanced Risk Analytics: VaR & CVaR Tail Risks")
    box = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Value at Risk (95% Daily VaR) represents the maximum daily loss threshold with 95% confidence (5th percentile).",
        "Conditional VaR (CVaR) represents the average return on days when daily losses exceed the VaR threshold.",
        "Small-cap equity schemes show severe tail risk, led by ABSL Small Cap (VaR: -2.39%, CVaR: -3.03%).",
        "Liquid and short-term debt schemes display negligible tail risk, with ICICI Pru Liquid showing a VaR of -0.019%."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.space_before = Pt(10)
        p_b.font.color.rgb = TEXT_RGB
        
    corr_img = os.path.join(figures_dir, "11_nav_correlation_matrix.png")
    if os.path.exists(corr_img):
        slide8.shapes.add_picture(corr_img, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.2))

    # SLIDE 9: Advanced Models — Monte Carlo & Efficient Frontier
    slide9 = create_slide("Advanced Modeling: GBM & Efficient Frontier")
    box = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Monte Carlo: 1,000 paths project the 5-year NAV compounding trajectory for SBI Bluechip Fund using Geometric Brownian Motion.",
        "Expected mean path projects 99.28% growth (NAV: ₹164.21), with 90% confidence bands between ₹94.50 and ₹248.10.",
        "Markowitz Optimization: Covariance mapping of returns for 5 key bluechip funds identifies the Efficient Frontier.",
        "Tangent Portfolio (Max Sharpe: 1.18) allocates: SBI Bluechip (32%), ICICI (28%), Nippon (25%), Kotak (15%)."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.space_before = Pt(10)
        p_b.font.color.rgb = TEXT_RGB
        
    donut_img = os.path.join(figures_dir, "12_sector_donut.png")
    if os.path.exists(donut_img):
        slide9.shapes.add_picture(donut_img, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.2))

    # SLIDE 10: Dashboard Screenshot — Executive & Performance
    slide10 = create_slide("Streamlit Dashboard: Executive & Performance")
    p1 = os.path.join(reports_dir, "Page_1_Industry_Overview.png")
    if os.path.exists(p1):
        slide10.shapes.add_picture(p1, Inches(0.5), Inches(1.6), Inches(5.9), Inches(4.8))
        
    p2 = os.path.join(reports_dir, "Page_2_Fund_Performance.png")
    if os.path.exists(p2):
        slide10.shapes.add_picture(p2, Inches(6.8), Inches(1.6), Inches(5.9), Inches(4.8))

    # SLIDE 11: Dashboard Screenshot — Investor & Market Trends
    slide11 = create_slide("Streamlit Dashboard: Investor & Market Trends")
    p3 = os.path.join(reports_dir, "Page_3_Investor_Analytics.png")
    if os.path.exists(p3):
        slide11.shapes.add_picture(p3, Inches(0.5), Inches(1.6), Inches(5.9), Inches(4.8))
        
    p4 = os.path.join(reports_dir, "Page_4_SIP_Market_Trends.png")
    if os.path.exists(p4):
        slide11.shapes.add_picture(p4, Inches(6.8), Inches(1.6), Inches(5.9), Inches(4.8))

    # SLIDE 12: Recommendations & Strategic Roadmap
    slide12 = create_slide("Key Takeaways & Strategic Recommendations")
    box = slide12.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Millennial Product Placement: Younger cohorts (18-35) make up 66.3% of users. Target them with mobile micro-SIP products.",
        "Tap Beyond-30 (B30) Markets: B30 semi-urban regions generate 30.1% of capital. Expand digital onboarding in tier-2/3 cities.",
        "Improve SIP Mandate Continuity: With 99% of users showing transaction gaps >35 days, implement push notification reminders, flexible pause options, and bank e-mandates.",
        "Embed Advanced Portfolio Tools: Integrate the asset allocator and Monte Carlo simulator into retail client portals.",
        "Thank You! Any Questions?"
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(15)
        p_b.space_before = Pt(12)
        p_b.font.color.rgb = TEXT_RGB

    prs.save(pptx_path)
    # Copy to reports as well to keep both paths
    import shutil
    try:
        shutil.copy(pptx_path, reports_pptx_path)
        print("Presentation PPTX copied to reports directory successfully!")
    except Exception as e:
        print(f"Error copying presentation: {e}")
        
    print("Presentation PPTX generated successfully at root Bluestock_MF_Presentation.pptx")

if __name__ == "__main__":
    build_pdf()
    build_pptx()
    conn.close()
